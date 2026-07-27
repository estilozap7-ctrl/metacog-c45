"""
=========================================================
PyC45 Framework — Generador de Presentaciones (PPTX & PDF)
=========================================================
Este script genera automáticamente dos formatos de presentación:
1. PowerPoint (.pptx) usando python-pptx.
2. PDF Paisaje (.pdf) usando fpdf2.

Ambas presentaciones tienen un diseño oscuro premium coherente
con el sistema de diseño visual de PyC45.
=========================================================
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from fpdf import FPDF

# =========================================================
# CONFIGURACIÓN DE COLORES DEL TEMA (PyC45 Dark Theme)
# =========================================================
COLOR_BG = (10, 11, 16)         # #0a0b10 (Fondo Principal)
COLOR_CARD = (21, 24, 34)       # #151822 (Fondo de Tarjetas)
COLOR_ACCENT = (108, 99, 255)   # #6c63ff (Acento Púrpura)
COLOR_ACCENT_LT = (139, 131, 255) # #8b83ff (Púrpura Claro)
COLOR_SUCCESS = (0, 212, 170)    # #00d4aa (Teal/Menta)
COLOR_WARNING = (255, 217, 61)   # #ffd93d (Amarillo)
COLOR_TEXT_MAIN = (228, 228, 231) # #e4e4e7 (Texto Principal)
COLOR_TEXT_MUTED = (161, 161, 170) # #a1a1aa (Texto Secundario)
COLOR_BORDER = (40, 42, 54)      # Borde sutil

# =========================================================
# FUNCIÓN PARA LIMPIAR TEXTO EXCLUSIVO PARA FPDF (PDF)
# =========================================================
def clean_pdf_text(text):
    """Sustituye emojis y caracteres matemáticos/gráficos no soportados en Latin-1 por texto legible."""
    if not text:
        return ""
    
    # Emojis comunes
    text = text.replace("📊", "").replace("🧠", "").replace("⚡", "")
    text = text.replace("⚙️", "").replace("⚙", "").replace("🌿", "")
    text = text.replace("🔍", "").replace("🎓", "").replace("📈", "")
    
    # Caracteres de árbol
    text = text.replace("├──", "|--").replace("└──", "`--")
    text = text.replace("│", "|").replace("─", "-").replace("├", "|")
    
    # Símbolos matemáticos e indicadores
    text = text.replace("Σ", "Sum").replace("•", "-")
    text = text.replace("➔", "->").replace("->", "->")
    text = text.replace("—", "-").replace("–", "-")
    
    # Filtrar cualquier caracter que no se pueda codificar en latin-1
    cleaned = []
    for char in text:
        try:
            char.encode('latin-1')
            cleaned.append(char)
        except UnicodeEncodeError:
            cleaned.append(' ')
            
    return "".join(cleaned)

# =========================================================
# AUXILIARES PARA POWERPOINT (PPTX)
# =========================================================
def set_pptx_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_rgb)

def create_pptx_base_slide(prs, title, slide_num):
    blank_layout = prs.slide_layouts[6] # blank layout
    slide = prs.slides.add_slide(blank_layout)
    set_pptx_slide_background(slide, COLOR_BG)
    
    # Agregar número de diapositiva y título
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    
    # Nro Diapositiva
    p_num = tf.paragraphs[0]
    p_num.text = f"DIAPOSITIVA {slide_num:02d}"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(10)
    p_num.font.bold = True
    p_num.font.color.rgb = RGBColor(*COLOR_ACCENT)
    
    # Título
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    p_title.space_before = Pt(4)
    
    return slide

def add_pptx_card(slide, left, top, width, height, title, text, bg_color=COLOR_CARD, border_color=COLOR_ACCENT, title_color=COLOR_ACCENT_LT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    # Título de tarjeta
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(15)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(*title_color)
    p_title.space_after = Pt(6)
    
    # Texto
    p_text = tf.add_paragraph()
    p_text.text = text
    p_text.font.name = "Arial"
    p_text.font.size = Pt(11)
    p_text.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p_text.line_spacing = 1.2
    
    return shape

# =========================================================
# AUXILIARES PARA PDF
# =========================================================
class LandscapePDF(FPDF):
    def __init__(self):
        super().__init__(orientation="L", unit="mm", format="A4") # 297mm x 210mm
        # Cargar fuentes del sistema
        font_path_regular = r"C:\Windows\Fonts\arial.ttf"
        font_path_bold = r"C:\Windows\Fonts\arialbd.ttf"
        font_path_italic = r"C:\Windows\Fonts\ariali.ttf"
        
        self.add_font("Arial", "", font_path_regular)
        self.add_font("Arial", "B", font_path_bold)
        self.add_font("Arial", "I", font_path_italic)

    def draw_bg(self):
        self.set_fill_color(*COLOR_BG)
        self.rect(0, 0, 297, 210, "F")
        # Decoraciones circulares de fondo
        self.set_fill_color(20, 20, 35)
        self.ellipse(-20, -20, 100, 100, "F")
        self.set_fill_color(15, 25, 30)
        self.ellipse(220, 140, 110, 110, "F")

    def create_base_slide(self, title, slide_num):
        self.add_page()
        self.draw_bg()
        
        # Nro Diapositiva
        self.set_xy(15, 10)
        self.set_font("Arial", "B", 9)
        self.set_text_color(*COLOR_ACCENT)
        cleaned_num = clean_pdf_text(f"DIAPOSITIVA {slide_num:02d}")
        self.cell(0, 5, cleaned_num, border=0, new_x="LMARGIN", new_y="NEXT")
        
        # Título
        self.set_font("Arial", "B", 22)
        self.set_text_color(255, 255, 255)
        cleaned_title = clean_pdf_text(title)
        self.cell(0, 10, cleaned_title, border=0, new_x="LMARGIN", new_y="NEXT")
        
        # Línea de cabecera
        self.set_draw_color(*COLOR_BORDER)
        self.set_line_width(0.3)
        self.line(15, 28, 282, 28)
        
        # Pie de página
        self.set_xy(15, 198)
        self.set_font("Arial", "I", 8)
        self.set_text_color(*COLOR_TEXT_MUTED)
        self.cell(100, 5, clean_pdf_text("PyC45: Framework de Clasificación C4.5"), border=0)
        self.set_x(260)
        self.cell(0, 5, clean_pdf_text(f"Pág. {slide_num} / 9"), border=0, align="R")

    def draw_card(self, x, y, w, h, title, text, bg_color=COLOR_CARD, border_color=COLOR_ACCENT, title_color=COLOR_ACCENT_LT):
        self.set_fill_color(*bg_color)
        self.set_draw_color(*border_color)
        self.set_line_width(0.4)
        self.rect(x, y, w, h, "FD")
        
        # Título Tarjeta
        self.set_xy(x + 5, y + 4)
        self.set_font("Arial", "B", 12)
        self.set_text_color(*title_color)
        self.cell(w - 10, 5, clean_pdf_text(title), border=0, new_x="LMARGIN", new_y="NEXT")
        
        # Texto Tarjeta
        self.set_xy(x + 5, y + 10)
        self.set_font("Arial", "", 9)
        self.set_text_color(*COLOR_TEXT_MAIN)
        self.multi_cell(w - 10, 4.2, clean_pdf_text(text), border=0)


def generate_presentations():
    # Rutas
    base_dir = os.path.dirname(os.path.abspath(__file__))
    pptx_path = os.path.join(base_dir, 'presentacion_pyc45.pptx')
    pdf_path = os.path.join(base_dir, 'presentacion_pyc45.pdf')
    
    print("Creando presentación PowerPoint (.pptx)...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5) # 16:9
    
    print("Creando presentación PDF (.pdf)...")
    pdf = LandscapePDF()
    
    # =========================================================
    # DIAPOSITIVA 1: PORTADA
    # =========================================================
    # PPTX
    slide_layout = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(slide_layout)
    set_pptx_slide_background(s1, COLOR_BG)
    
    # Decoración visual en portada
    shape = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5), Inches(6), Inches(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(18, 19, 28)
    shape.line.fill.background()
    
    # Título e Información en Portada
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_badge = tf.paragraphs[0]
    p_badge.text = "📊 FRAMEWORK PyC45"
    p_badge.font.name = "Arial"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = RGBColor(*COLOR_ACCENT)
    p_badge.space_after = Pt(20)
    
    p_title = tf.add_paragraph()
    p_title.text = "PyC45: Clasificación C4.5 en Python"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Funcionamiento paso a paso, características y aportes a la analítica de datos explicable"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(*COLOR_TEXT_MUTED)
    p_sub.space_before = Pt(8)
    p_sub.space_after = Pt(50)
    
    p_meta = tf.add_paragraph()
    p_meta.text = "Autor: Luis Alberto Buelvas Cogollo   |   Caso Práctico: Scoring Crediticio (UCI)   |   PyC45 v1.0.0"
    p_meta.font.name = "Arial"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = RGBColor(*COLOR_ACCENT_LT)
    
    # PDF Slide 1
    pdf.add_page()
    pdf.draw_bg()
    pdf.set_xy(20, 50)
    pdf.set_font("Arial", "B", 12)
    pdf.set_text_color(*COLOR_ACCENT)
    pdf.cell(0, 5, clean_pdf_text("FRAMEWORK PyC45"), border=0, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)
    pdf.set_font("Arial", "B", 32)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 15, clean_pdf_text("PyC45: Clasificación C4.5 en Python"), border=0, new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Arial", "", 16)
    pdf.set_text_color(*COLOR_TEXT_MUTED)
    pdf.cell(0, 10, clean_pdf_text("Funcionamiento paso a paso, características y aportes a la analítica de datos explicable"), border=0, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(35)
    pdf.set_font("Arial", "B", 10)
    pdf.set_text_color(*COLOR_ACCENT_LT)
    pdf.cell(0, 5, clean_pdf_text("Autor: Luis Alberto Buelvas Cogollo    |    Caso Práctico: Scoring Crediticio (UCI)    |    PyC45 v1.0.0"), border=0, align="L")

    # =========================================================
    # DIAPOSITIVA 2: ¿QUÉ ES PYC45? (OBJETIVO Y FUNCIONALIDAD)
    # =========================================================
    # PPTX
    s2 = create_pptx_base_slide(prs, "¿Qué es el Proyecto PyC45?", 2)
    
    # Columna Izquierda: Descripción
    left_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5))
    tf2 = left_box.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "PyC45 es una implementación limpia, modular y orientada a objetos en Python puro del algoritmo clásico C4.5 diseñado por Ross Quinlan en 1993."
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p.space_after = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "Objetivo y Funcionalidad:"
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLOR_ACCENT_LT)
    p.space_after = Pt(8)
    
    p = tf2.add_paragraph()
    p.text = "• Clasificación Binaria: Resuelve problemas de aprendizaje supervisado con salida 0 o 1 (ej. abandono de clientes, fraude bancario o scoring de crédito).\n" \
             "• Transparencia ('Caja Blanca'): A diferencia de redes neuronales, genera reglas explícitas de decisión auditables y fáciles de comprender.\n" \
             "• Modularidad: Facilita la personalización matemática (fórmulas de pureza y criterios de parada) de manera sencilla."
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MUTED)
    p.line_spacing = 1.3
    
    # Columna Derecha: Tarjetas
    add_pptx_card(s2, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.1), 
                  "🧠 Propósito Educativo", 
                  "Código transparente en español estructurado sin optimizaciones complejas de bajo nivel. Diseñado para que estudiantes y desarrolladores entiendan el comportamiento íntimo del algoritmo y su flujo de datos recursivo.")
    
    add_pptx_card(s2, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.1), 
                  "⚡ Enfoque Práctico", 
                  "Preparado para aplicarse a conjuntos de datos comerciales del mundo real. Cuenta con soporte nativo de variables numéricas y el algoritmo de Poda por Error Reducido (REP) para maximizar la generalización.")

    # PDF Slide 2
    pdf.create_base_slide("¿Qué es el Proyecto PyC45?", 2)
    # Columna Izquierda
    pdf.set_xy(15, 35)
    pdf.set_font("Arial", "", 12)
    pdf.set_text_color(*COLOR_TEXT_MAIN)
    pdf.multi_cell(125, 6, clean_pdf_text("PyC45 es una implementación limpia, modular y orientada a objetos en Python puro del algoritmo clásico C4.5 diseñado por Ross Quinlan en 1993."), border=0)
    pdf.ln(5)
    pdf.set_font("Arial", "B", 11)
    pdf.set_text_color(*COLOR_ACCENT_LT)
    pdf.cell(0, 5, clean_pdf_text("Objetivo y Funcionalidad:"), border=0, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.set_font("Arial", "", 10)
    pdf.set_text_color(*COLOR_TEXT_MUTED)
    pdf.multi_cell(125, 5, clean_pdf_text("• Clasificación Binaria: Resuelve problemas de aprendizaje supervisado con salida 0 o 1 (ej. abandono de clientes, fraude bancario o scoring de crédito).\n• Transparencia ('Caja Blanca'): A diferencia de redes neuronales, genera reglas explícitas de decisión auditables y fáciles de comprender.\n• Modularidad: Facilita la personalización matemática de manera sencilla."), border=0)
    # Columna Derecha: Tarjetas
    pdf.draw_card(152, 35, 130, 48, "🧠 Propósito Educativo", "Código transparente en español estructurado sin optimizaciones complejas de bajo nivel. Diseñado para que estudiantes y desarrolladores entiendan el comportamiento íntimo del algoritmo y su flujo de datos recursivo.")
    pdf.draw_card(152, 90, 130, 48, "⚡ Enfoque Práctico", "Preparado para aplicarse a conjuntos de datos comerciales del mundo real. Cuenta con soporte nativo de variables numéricas y el algoritmo de Poda por Error Reducido (REP) para maximizar la generalización.")

    # =========================================================
    # DIAPOSITIVA 3: ARQUITECTURA MODULAR
    # =========================================================
    # PPTX
    s3 = create_pptx_base_slide(prs, "Arquitectura y Estructura Modular", 3)
    
    # Estructura del repositorio en caja de código (Izquierda)
    code_box = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(15, 16, 23)
    code_box.line.color.rgb = RGBColor(*COLOR_BORDER)
    
    tf_code = code_box.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = Inches(0.3)
    tf_code.margin_top = Inches(0.3)
    p_c = tf_code.paragraphs[0]
    p_c.text = "Estructura del Proyecto PyC45:\n\n" \
               "PyC45/\n" \
               "├── core/                    # Núcleo Algorítmico\n" \
               "│   ├── classifier.py        # fit / predict / score\n" \
               "│   ├── tree_builder.py      # Construcción recursiva\n" \
               "│   ├── best_split.py        # Puntos de corte continuos\n" \
               "│   ├── entropy.py           # Shannon Entropy\n" \
               "│   └── pruning.py           # Poda REP\n" \
               "├── metrics/                 # Suite de Validación\n" \
               "│   ├── confusion_matrix.py  # Estructura de aciertos\n" \
               "│   ├── classification_metrics.py # MCC, F1, Exactitud\n" \
               "│   └── roc_auc.py           # Curva ROC y AUC\n" \
               "└── visualization/           # Explicabilidad y UI\n" \
               "    ├── tree_plot.py         # Gráfico con Matplotlib\n" \
               "    └── feature_importance.py# Gain Ratio acumulado"
    p_c.font.name = "Courier New"
    p_c.font.size = Pt(10)
    p_c.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p_c.line_spacing = 1.15
    
    # Tarjetas Derecha
    add_pptx_card(s3, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.1), 
                  "⚙️ Core Algorítmico Desacoplado", 
                  "La construcción del árbol se separa de la matemática de pureza e información. El clasificador interactúa con un generador recursivo (TreeBuilder), logrando un diseño de bajo acoplamiento y alta cohesión.")
    
    add_pptx_card(s3, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.1), 
                  "📊 Evaluación Multidimensional", 
                  "Módulos nativos e independientes para la evaluación del desempeño. Integra métricas robustas frente a desbalances como el Coeficiente de Matthews (MCC) y la Curva ROC-AUC sin dependencias externas.")

    # PDF Slide 3
    pdf.create_base_slide("Arquitectura y Estructura Modular", 3)
    # Caja de código Izquierda
    pdf.set_fill_color(15, 16, 23)
    pdf.set_draw_color(*COLOR_BORDER)
    pdf.rect(15, 35, 125, 140, "FD")
    pdf.set_xy(18, 38)
    pdf.set_font("Courier", "", 8)
    pdf.set_text_color(*COLOR_TEXT_MAIN)
    
    # Generar el árbol limpio para PDF
    pdf_code_tree = "Estructura del Proyecto PyC45:\n\n" \
                    "PyC45/\n" \
                    "|-- core/                    # Nucleo Algoritmico\n" \
                    "|   |-- classifier.py        # fit / predict / score\n" \
                    "|   |-- tree_builder.py      # Construccion recursiva\n" \
                    "|   |-- best_split.py        # Cortes continuos\n" \
                    "|   |-- entropy.py           # Shannon Entropy\n" \
                    "|   `-- pruning.py           # Poda REP\n" \
                    "|-- metrics/                 # Suite de Validacion\n" \
                    "|   |-- confusion_matrix.py\n" \
                    "|   |-- classification_metrics.py\n" \
                    "|   `-- roc_auc.py           # Curva ROC y AUC\n" \
                    "`-- visualization/           # Explicabilidad y UI\n" \
                    "    |-- tree_plot.py         # Matplotlib Plot\n" \
                    "    `-- feature_importance.py"
    
    pdf.multi_cell(120, 4.3, clean_pdf_text(pdf_code_tree), border=0)
    # Tarjetas Derecha
    pdf.draw_card(152, 35, 130, 48, "⚙️ Core Algorítmico Desacoplado", "La construcción del árbol se separa de la matemática de pureza e información. El clasificador interactúa con un generador recursivo (TreeBuilder), logrando un diseño de bajo acoplamiento y alta cohesión.")
    pdf.draw_card(152, 90, 130, 48, "📊 Evaluación Multidimensional", "Módulos nativos e independientes para la evaluación del desempeño. Integra métricas robustas frente a desbalances como el Coeficiente de Matthews (MCC) y la Curva ROC-AUC sin dependencias externas.")

    # =========================================================
    # DIAPOSITIVA 4: FUNCIONAMIENTO PASO A PASO (EL PIPELINE)
    # =========================================================
    # PPTX
    s4 = create_pptx_base_slide(prs, "Funcionamiento Paso a Paso (El Pipeline)", 4)
    
    # Diagrama de flujo de nodos (rectángulos encadenados)
    steps = [
        ("1. Carga de Datos", Inches(0.8), COLOR_CARD),
        ("2. Preproceso", Inches(3.8), COLOR_CARD),
        ("3. Partición (70/30)", Inches(6.8), COLOR_CARD),
        ("4. Creador C4.5", Inches(9.8), COLOR_CARD),
    ]
    for name, left, color in steps:
        shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.7), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(*COLOR_ACCENT)
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
    steps_row2 = [
        ("7. Visualización", Inches(0.8), COLOR_CARD),
        ("6. Evaluación", Inches(3.8), COLOR_CARD),
        ("5. Poda REP", Inches(6.8), COLOR_CARD),
    ]
    for name, left, color in steps_row2:
        shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.0), Inches(2.7), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(*COLOR_ACCENT)
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
    # Tarjetas informativas de las fases en la parte inferior
    add_pptx_card(s4, Inches(0.8), Inches(4.3), Inches(3.7), Inches(2.2), 
                  "Fase A: Preparación", 
                  "Carga automática desde Pandas. Remoción de columnas nulas e irrelevantes (ej. ID). Partición robusta entrenamiento/validación (70/30).", 
                  border_color=COLOR_BORDER)
                  
    add_pptx_card(s4, Inches(4.8), Inches(4.3), Inches(3.7), Inches(2.2), 
                  "Fase B: Modelado", 
                  "Entrenamiento recursivo top-down basado en la maximización del Gain Ratio. Aplicación bottom-up de la poda REP para erradicar sobreajuste.", 
                  border_color=COLOR_BORDER)
                  
    add_pptx_card(s4, Inches(8.8), Inches(4.3), Inches(3.7), Inches(2.2), 
                  "Fase C: Análisis", 
                  "Validación final con métricas robustas. Generación del gráfico de árbol con Matplotlib, ranking de variables e informe web dinámico.", 
                  border_color=COLOR_BORDER)

    # PDF Slide 4
    pdf.create_base_slide("Funcionamiento Paso a Paso (El Pipeline)", 4)
    # Fila 1 diagramas
    pdf.set_fill_color(*COLOR_CARD)
    pdf.set_draw_color(*COLOR_ACCENT)
    pdf.set_line_width(0.4)
    for i, step_name in enumerate(["1. Carga Datos", "2. Preproceso", "3. Particion (70/30)", "4. Creador C4.5"]):
        pdf.rect(15 + i*68, 35, 60, 15, "FD")
        pdf.set_xy(15 + i*68, 40)
        pdf.set_font("Arial", "B", 10)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(60, 5, clean_pdf_text(step_name), border=0, align="C")
        if i < 3:
            pdf.set_xy(75 + i*68, 40)
            pdf.set_font("Arial", "B", 12)
            pdf.set_text_color(*COLOR_ACCENT_LT)
            pdf.cell(8, 5, "->", border=0, align="C")
            
    # Fila 2 diagramas
    for i, step_name in enumerate(["5. Poda REP", "6. Evaluacion", "7. Visualizacion"]):
        pdf.rect(219 - i*68, 60, 60, 15, "FD")
        pdf.set_xy(219 - i*68, 65)
        pdf.set_font("Arial", "B", 10)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(60, 5, clean_pdf_text(step_name), border=0, align="C")
        if i < 2:
            pdf.set_xy(211 - i*68, 65)
            pdf.set_font("Arial", "B", 12)
            pdf.set_text_color(*COLOR_ACCENT_LT)
            pdf.cell(8, 5, "<-", border=0, align="C")
            
    # Tarjetas fase inferior
    pdf.draw_card(15, 90, 84, 52, "Fase A: Preparación", "Carga automática desde Pandas. Remoción de columnas nulas e irrelevantes (ej. ID). Partición robusta entrenamiento/validación (70/30).", border_color=COLOR_BORDER)
    pdf.draw_card(106.5, 90, 84, 52, "Fase B: Modelado", "Entrenamiento recursivo top-down basado en la maximización del Gain Ratio. Aplicación bottom-up de la poda REP para erradicar sobreajuste.", border_color=COLOR_BORDER)
    pdf.draw_card(198, 90, 84, 52, "Fase C: Análisis", "Validación final con métricas robustas. Generación del gráfico de árbol con Matplotlib, ranking de variables e informe web dinámico.", border_color=COLOR_BORDER)

    # =========================================================
    # DIAPOSITIVA 5: EL NÚCLEO MATEMÁTICO: CRITERIO C4.5
    # =========================================================
    # PPTX
    s5 = create_pptx_base_slide(prs, "El Núcleo Matemático: Criterio C4.5", 5)
    
    # Tarjetas Fórmulas Izquierda
    add_pptx_card(s5, Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.1), 
                  "1. Entropía de Shannon (Pureza)", 
                  "H(S) = - Σ p_i log2(p_i)\n" \
                  "Mide la impureza o incertidumbre del conjunto de clases S en un nodo. Varía entre 0.0 (máxima pureza) y 1.0 (clases equipartidas).",
                  title_color=COLOR_SUCCESS)
                  
    add_pptx_card(s5, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.1), 
                  "2. Ganancia de Información", 
                  "IG(S, A) = H(S) - Σ (|S_v| / |S|) * H(S_v)\n" \
                  "Mide la reducción neta de entropía al dividir por la variable A. Presenta un fuerte sesgo a atributos con alta cardinalidad.",
                  title_color=COLOR_SUCCESS)
                  
    # Tarjeta Gain Ratio Derecha
    add_pptx_card(s5, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), 
                  "3. Gain Ratio (Normalización C4.5)", 
                  "El algoritmo C4.5 resuelve el sesgo del Gain de Información normalizándolo por el Split Information:\n\n" \
                  "• SplitInfo(S, A) = - Σ (|S_v| / |S|) * log2(|S_v| / |S|)\n" \
                  "• GainRatio(S, A) = IG(S, A) / SplitInfo(S, A)\n\n" \
                  "El Split Information mide la entropía inherente de la partición (cuántos subconjuntos se generan). Si una variable toma un valor único para cada registro (ej. ID), el SplitInfo será máximo, reduciendo drásticamente su Gain Ratio y evitando divisiones artificiales.",
                  border_color=COLOR_SUCCESS)

    # PDF Slide 5
    pdf.create_base_slide("El Núcleo Matemático: Criterio C4.5", 5)
    # Tarjetas Izquierda
    pdf.draw_card(15, 35, 125, 48, "1. Entropía de Shannon (Pureza)", "H(S) = - Sum p_i log2(p_i)\nMide la impureza o incertidumbre del conjunto de clases S en un nodo. Varía entre 0.0 (máxima pureza) y 1.0 (clases equipartidas).", title_color=COLOR_SUCCESS)
    pdf.draw_card(15, 90, 125, 48, "2. Ganancia de Información", "IG(S, A) = H(S) - Sum (|S_v| / |S|) * H(S_v)\nMide la reducción neta de entropía al dividir por la variable A. Presenta un fuerte sesgo a atributos con alta cardinalidad.", title_color=COLOR_SUCCESS)
    # Tarjeta Derecha
    pdf.draw_card(152, 35, 130, 103, "3. Gain Ratio (Normalización C4.5)", "El algoritmo C4.5 resuelve el sesgo del Gain de Información normalizándolo por el Split Information:\n\n• SplitInfo(S, A) = - Sum (|S_v| / |S|) * log2(|S_v| / |S|)\n• GainRatio(S, A) = IG(S, A) / SplitInfo(S, A)\n\nEl Split Information mide la entropía inherente de la partición (cuántos subconjuntos se generan). Si una variable toma un valor único para cada registro (ej. ID), el SplitInfo será máximo, reduciendo drásticamente su Gain Ratio y evitando divisiones artificiales.", border_color=COLOR_SUCCESS)

    # =========================================================
    # DIAPOSITIVA 6: SIMPLIFICACIÓN Y CONTROL: PODA REP
    # =========================================================
    # PPTX
    s6 = create_pptx_base_slide(prs, "Simplificación y Control: Poda REP", 6)
    
    left_box6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5))
    tf6 = left_box6.text_frame
    tf6.word_wrap = True
    
    p = tf6.paragraphs[0]
    p.text = "La Poda por Error Reducido (Reduced Error Pruning - REP) es una técnica de simplificación estructurada del árbol posterior al entrenamiento."
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p.space_after = Pt(20)
    
    p = tf6.add_paragraph()
    p.text = "Características Clave de la Poda:"
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLOR_ACCENT_LT)
    p.space_after = Pt(8)
    
    p = tf6.add_paragraph()
    p.text = "• Enfoque Bottom-Up: Recorre el árbol de abajo hacia arriba comenzando por las ramas más profundas.\n" \
             "• Basado en Validación: Utiliza un conjunto de datos de validación independiente (no visto en el entrenamiento) para medir el error real.\n" \
             "• Colapso Inteligente: Si al consolidar un subárbol en una hoja (que predice la clase mayoritaria) la exactitud local no empeora, la poda es ejecutada."
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MUTED)
    p.line_spacing = 1.3
    
    # Tarjetas Derecha
    add_pptx_card(s6, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.1), 
                  "Antes de la Poda (Sobreajuste)", 
                  "El árbol crece memorizando particularidades del conjunto de entrenamiento (incluyendo ruido). Posee alta profundidad, gran cantidad de hojas, pero menor exactitud (accuracy) sobre datos de prueba.",
                  border_color=COLOR_WARNING, title_color=COLOR_WARNING)
                  
    add_pptx_card(s6, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.1), 
                  "Después de la Poda (Generalización)", 
                  "Se eliminan las ramas con nulo poder predictivo. Esto mejora la generalización en datos futuros, reduce el tamaño del árbol y aumenta la comprensibilidad del modelo final.",
                  border_color=COLOR_SUCCESS, title_color=COLOR_SUCCESS)

    # PDF Slide 6
    pdf.create_base_slide("Simplificación y Control: Poda REP", 6)
    # Izquierda
    pdf.set_xy(15, 35)
    pdf.set_font("Arial", "", 12)
    pdf.set_text_color(*COLOR_TEXT_MAIN)
    pdf.multi_cell(125, 6, clean_pdf_text("La Poda por Error Reducido (Reduced Error Pruning - REP) es una técnica de simplificación estructurada del árbol posterior al entrenamiento."), border=0)
    pdf.ln(5)
    pdf.set_font("Arial", "B", 11)
    pdf.set_text_color(*COLOR_ACCENT_LT)
    pdf.cell(0, 5, clean_pdf_text("Características Clave de la Poda:"), border=0, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.set_font("Arial", "", 10)
    pdf.set_text_color(*COLOR_TEXT_MUTED)
    pdf.multi_cell(125, 5, clean_pdf_text("• Enfoque Bottom-Up: Recorre el árbol de abajo hacia arriba comenzando por las ramas más profundas.\n• Basado en Validación: Utiliza un conjunto de datos de validación independiente (no visto en el entrenamiento) para medir el error real.\n• Colapso Inteligente: Si al consolidar un subárbol en una hoja la exactitud local no empeora, la poda es ejecutada."), border=0)
    # Derecha
    pdf.draw_card(152, 35, 130, 48, "Antes de la Poda (Sobreajuste)", "El árbol crece memorizando particularidades del conjunto de entrenamiento (incluyendo ruido). Posee alta profundidad, gran cantidad de hojas, pero menor exactitud (accuracy) sobre datos de prueba.", border_color=COLOR_WARNING, title_color=COLOR_WARNING)
    pdf.draw_card(152, 90, 130, 48, "Después de la Poda (Generalización)", "Se eliminan las ramas con nulo poder predictivo. Esto mejora la generalización en datos futuros, reduce el tamaño del árbol y aumenta la comprensibilidad del modelo final.", border_color=COLOR_SUCCESS, title_color=COLOR_SUCCESS)

    # =========================================================
    # DIAPOSITIVA 7: PYC45 VS ALGORITMO C4.5 CLÁSICO
    # =========================================================
    # PPTX
    s7 = create_pptx_base_slide(prs, "PyC45 vs. Algoritmo C4.5 Clásico", 7)
    
    # Crear Tabla comparativa nativa
    rows, cols = 6, 3
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
    table_shape = s7.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Encabezados
    headers = ["Característica", "C4.5 Clásico (Quinlan, 1993)", "Framework PyC45"]
    for i, h_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(18, 19, 28)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(*COLOR_ACCENT_LT)
        
    # Datos de la tabla
    row_data = [
        ["Lenguaje", "Código en C monolítico estructurado", "Python 3 modular orientado a objetos"],
        ["Estructuras", "Variables directas con memoria manual", "Estructura de nodos (DecisionNode) encapsulada"],
        ["Formatos de Datos", "Archivos planos propietarios (.data, .names)", "Integración nativa con Pandas (DataFrames) y NumPy"],
        ["Estrategia de Poda", "Poda pesimista basada en cotas teóricas", "Poda por Error Reducido (REP) con validación"],
        ["Métricas e Informes", "Tablas de confusión y error sencillas en consola", "MCC, F1, Curva ROC-AUC, Importancia, Reporte Web"]
    ]
    
    for r_idx, row_values in enumerate(row_data):
        for c_idx, val in enumerate(row_values):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*COLOR_CARD)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
            
    # PDF Slide 7
    pdf.create_base_slide("PyC45 vs. Algoritmo C4.5 Clásico", 7)
    # Dibujar tabla
    col_widths = [45, 110, 110]
    row_height = 20
    start_x = 15
    start_y = 35
    
    # Encabezados PDF
    pdf.set_font("Arial", "B", 10)
    pdf.set_fill_color(18, 19, 28)
    pdf.set_text_color(*COLOR_ACCENT_LT)
    
    pdf.set_xy(start_x, start_y)
    pdf.cell(col_widths[0], 10, clean_pdf_text("Caracteristica"), border=1, fill=True, align="L")
    pdf.cell(col_widths[1], 10, clean_pdf_text("C4.5 Clasico (Quinlan, 1993)"), border=1, fill=True, align="C")
    pdf.cell(col_widths[2], 10, clean_pdf_text("Framework PyC45"), border=1, fill=True, align="C")
    
    # Filas PDF
    pdf.set_font("Arial", "", 9)
    pdf.set_fill_color(*COLOR_CARD)
    pdf.set_text_color(*COLOR_TEXT_MAIN)
    
    for r_idx, row_values in enumerate(row_data):
        pdf.set_xy(start_x, start_y + 10 + r_idx*row_height)
        pdf.cell(col_widths[0], row_height, clean_pdf_text(row_values[0]), border=1, fill=True, align="L")
        
        x_before = pdf.get_x()
        y_before = pdf.get_y()
        pdf.multi_cell(col_widths[1], row_height/2, clean_pdf_text(row_values[1]), border=1, fill=True, align="C")
        pdf.set_xy(x_before + col_widths[1], y_before)
        pdf.multi_cell(col_widths[2], row_height/2, clean_pdf_text(row_values[2]), border=1, fill=True, align="C")

    # =========================================================
    # DIAPOSITIVA 8: APORTE A LA ANALÍTICA DE DATOS
    # =========================================================
    # PPTX
    s8 = create_pptx_base_slide(prs, "Aporte a la Analítica de Datos", 8)
    
    # Tres tarjetas verticales
    add_pptx_card(s8, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.5), 
                  "🔍 Explicabilidad (XAI)", 
                  "En sectores fuertemente regulados como finanzas y salud, la explicabilidad es obligatoria.\n\nPyC45 entrega un conjunto explícito y comprensible de reglas de decisión (de tipo SI-ENTONCES) que pueden ser auditadas, validadas y explicadas directamente a las partes interesadas sin áreas grises ni dudas.",
                  border_color=COLOR_ACCENT)
                  
    add_pptx_card(s8, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.5), 
                  "🎓 Valor Pedagógico", 
                  "La mayoría de librerías modernas de machine learning ocultan sus algoritmos en complejos bindings de C/C++.\n\nPyC45 expone cada línea de código en Python limpio. Permite descomponer la ganancia de información, la selección de umbrales numéricos y el proceso recursivo de poda REP paso a paso.",
                  border_color=COLOR_ACCENT)
                  
    add_pptx_card(s8, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5), 
                  "📈 Métricas Realistas", 
                  "Evita el autoengaño frente a datasets desbalanceados de la vida real.\n\nAl integrar el Coeficiente de Matthews (MCC) y la Curva ROC-AUC de forma nativa, PyC45 provee al analista una perspectiva del rendimiento de clasificación mucho más rigurosa que la simple Exactitud global.",
                  border_color=COLOR_ACCENT)

    # PDF Slide 8
    pdf.create_base_slide("Aporte a la Analítica de Datos", 8)
    pdf.draw_card(15, 35, 84, 140, "🔍 Explicabilidad (XAI)", "En sectores fuertemente regulados como finanzas y salud, la explicabilidad es obligatoria.\n\nPyC45 entrega un conjunto explícito y comprensible de reglas de decisión (de tipo SI-ENTONCES) que pueden ser auditadas, validadas y explicadas directamente a las partes interesadas sin áreas grises ni dudas.", border_color=COLOR_ACCENT)
    pdf.draw_card(106.5, 35, 84, 140, "🎓 Valor Pedagógico", "La mayoría de librerías modernas de machine learning ocultan sus algoritmos en complejos bindings de C/C++.\n\nPyC45 expone cada línea de código en Python limpio. Permite descomponer la ganancia de información, la selección de umbrales numéricos y el proceso recursivo de poda REP paso a paso.", border_color=COLOR_ACCENT)
    pdf.draw_card(198, 35, 84, 140, "📈 Métricas Realistas", "Evita el autoengaño frente a datasets desbalanceados de la vida real.\n\nAl integrar el Coeficiente de Matthews (MCC) y la Curva ROC-AUC de forma nativa, PyC45 provee al analista una perspectiva del rendimiento de clasificación mucho más rigurosa que la simple Exactitud global.", border_color=COLOR_ACCENT)

    # =========================================================
    # DIAPOSITIVA 9: CONCLUSIONES Y DEMOSTRACIÓN
    # =========================================================
    # PPTX
    s9 = create_pptx_base_slide(prs, "Conclusiones del Proyecto", 9)
    
    # Izquierda
    left_box9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5))
    tf9 = left_box9.text_frame
    tf9.word_wrap = True
    
    p = tf9.paragraphs[0]
    p.text = "• Control y Autonomía: PyC45 demuestra que es posible implementar un pipeline de clasificación riguroso y transparente en pocas líneas de código, sin dependencias complejas."
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p.space_after = Pt(16)
    
    p = tf9.add_paragraph()
    p.text = "• Revelaciones de Negocio: En la prueba de default de crédito, el modelo descubrió inmediatamente que el estado de pagos más reciente (PAY_0) concentra más del 80% de la importancia predictiva."
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p.space_after = Pt(16)
    
    p = tf9.add_paragraph()
    p.text = "• Sólida Línea Base: Funciona como un modelo explicable inicial excelente (baseline) antes de implementar algoritmos más opacos de ensamble."
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    
    # Derecha: Consola comandos
    console_box = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
    console_box.fill.solid()
    console_box.fill.fore_color.rgb = RGBColor(15, 16, 23)
    console_box.line.color.rgb = RGBColor(*COLOR_ACCENT)
    console_box.line.width = Pt(1.5)
    
    tf_c = console_box.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.4)
    tf_c.margin_top = Inches(0.4)
    
    p_ct = tf_c.paragraphs[0]
    p_ct.text = "⚡ Comandos para comenzar en consola:\n\n" \
                "# 1. Instalar requerimientos\n" \
                "pip install -r requirements.txt\n\n" \
                "# 2. Correr pipeline principal sintético\n" \
                "python main.py\n\n" \
                "# 3. Probar con selector interactivo\n" \
                "python ejecutar_con_dataset.py\n\n" \
                "# 4. Generar datos del reporte web\n" \
                "python web_verification/generate_report.py\n\n" \
                "Consulte el README.md en el raíz para detalles."
    p_ct.font.name = "Courier New"
    p_ct.font.size = Pt(11)
    p_ct.font.color.rgb = RGBColor(*COLOR_TEXT_MAIN)
    p_ct.line_spacing = 1.25

    # PDF Slide 9
    pdf.create_base_slide("Conclusiones del Proyecto", 9)
    # Izquierda
    pdf.set_xy(15, 35)
    pdf.set_font("Arial", "", 11)
    pdf.set_text_color(*COLOR_TEXT_MAIN)
    pdf.multi_cell(125, 7, clean_pdf_text("• Control y Autonomía: PyC45 demuestra que es posible implementar un pipeline de clasificación riguroso y transparente en pocas líneas de código, sin dependencias complejas.\n\n• Revelaciones de Negocio: En la prueba de default de crédito, el modelo descubrió inmediatamente que el estado de pagos más reciente (PAY_0) concentra más del 80% de la importancia predictiva.\n\n• Sólida Línea Base: Funciona como un modelo explicable inicial excelente (baseline) antes de implementar algoritmos más opacos de ensamble."), border=0)
    # Derecha: Consola
    pdf.set_fill_color(15, 16, 23)
    pdf.set_draw_color(*COLOR_ACCENT)
    pdf.rect(152, 35, 130, 140, "FD")
    pdf.set_xy(156, 39)
    pdf.set_font("Arial", "", 8.5) # Usamos Arial en vez de Courier para evitar warnings
    pdf.set_text_color(*COLOR_TEXT_MAIN)
    cmd_text = "⚡ Comandos para comenzar en consola:\n\n" \
               "# 1. Instalar requerimientos\n" \
               "pip install -r requirements.txt\n\n" \
               "# 2. Correr pipeline principal sintetico\n" \
               "python main.py\n\n" \
               "# 3. Probar con selector interactivo\n" \
               "python ejecutar_con_dataset.py\n\n" \
               "# 4. Generar datos del reporte web\n" \
               "python web_verification/generate_report.py\n\n" \
               "Consulte el README.md en el raiz para detalles."
    pdf.multi_cell(122, 5, clean_pdf_text(cmd_text), border=0)

    # Save PPTX
    print(f"Guardando PowerPoint en: {pptx_path}")
    prs.save(pptx_path)
    
    # Save PDF
   # print(f"Guardando PDF en: {pdf_path}")
    # pdf.output(pdf_path)
    
    print("¡Ambas presentaciones se generaron con éxito!")


if __name__ == "__main__":
    generate_presentations()
